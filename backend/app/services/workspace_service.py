"""
Workspace Service — folder-based file manager with multi-format upload and folder-aware chat.
Supports: PDF, XLSX, PPTX, DOCX, TXT, CSV, MD
"""

import os
import logging
from datetime import datetime, timezone
from bson import ObjectId
from motor.motor_asyncio import AsyncIOMotorDatabase
from app.config import settings
from app.services.llm_service import llm_service

logger = logging.getLogger(__name__)

ALLOWED_EXTENSIONS = {".pdf", ".xlsx", ".xls", ".pptx", ".docx", ".txt", ".csv", ".md"}
MAX_TEXT_CHARS = 200_000
ACCESS_LEVELS = ("workspace", "shared", "private")


# ═══════════════════════════════════════════════════════════
# ACCESS CONTROL
# ═══════════════════════════════════════════════════════════


def _can_access(item: dict, requesting_user_id: str, owner_user_id: str) -> bool:
    """Check if requesting_user_id can access an item."""
    if requesting_user_id == owner_user_id:
        return True
    level = item.get("access_level", "workspace")
    if level == "workspace":
        return True
    if level == "shared":
        return requesting_user_id in (item.get("shared_with") or [])
    return False  # private


async def _check_folder_chain_access(db, folder_id: str, requesting_user_id: str, owner_user_id: str) -> bool:
    """Walk up the folder tree. If any ancestor is restricted, check access."""
    current_id = folder_id
    seen = set()
    while current_id and current_id not in seen:
        seen.add(current_id)
        folder = await db.workspace_folders.find_one({"_id": ObjectId(current_id)})
        if not folder:
            break
        if not _can_access(folder, requesting_user_id, owner_user_id):
            return False
        current_id = folder.get("parent_id")
    return True


async def set_access_level(
    db: AsyncIOMotorDatabase, user_id: str, item_type: str, item_id: str,
    access_level: str, shared_with: list[str] = None,
) -> dict | None:
    """Set access level on a file or folder. Only the owner can do this."""
    if access_level not in ACCESS_LEVELS:
        return None
    collection = db.workspace_folders if item_type == "folder" else db.workspace_files
    item = await collection.find_one({"_id": ObjectId(item_id), "user_id": user_id})
    if not item:
        return None
    update = {
        "access_level": access_level,
        "shared_with": shared_with or [],
        "locked_by": user_id if access_level != "workspace" else None,
    }
    await collection.update_one({"_id": ObjectId(item_id)}, {"$set": update})
    return update


async def get_access_info(db: AsyncIOMotorDatabase, user_id: str, item_type: str, item_id: str) -> dict | None:
    """Get current access settings for an item."""
    collection = db.workspace_folders if item_type == "folder" else db.workspace_files
    item = await collection.find_one(
        {"_id": ObjectId(item_id), "user_id": user_id},
        {"access_level": 1, "shared_with": 1, "locked_by": 1},
    )
    if not item:
        return None
    return {
        "access_level": item.get("access_level", "workspace"),
        "shared_with": item.get("shared_with", []),
        "locked_by": item.get("locked_by"),
    }


def _detect_file_type(filename: str) -> str | None:
    ext = os.path.splitext(filename)[1].lower()
    return {
        ".pdf": "pdf", ".xlsx": "xlsx", ".xls": "xlsx",
        ".pptx": "pptx", ".docx": "docx",
        ".txt": "txt", ".csv": "txt", ".md": "txt",
    }.get(ext)


def _extract_tables_from_pdf(file_path: str) -> list[dict]:
    """Extract structured tables from PDF using pdfplumber."""
    tables = []
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages[:50], 1):  # Limit to 50 pages
                page_tables = page.extract_tables()
                for t_idx, table in enumerate(page_tables):
                    if not table or len(table) < 2:
                        continue
                    # First row as headers
                    headers = [str(c or "").strip() for c in table[0]]
                    rows = []
                    for row in table[1:]:
                        cells = [str(c or "").strip() for c in row]
                        if any(cells):
                            rows.append(cells)
                    if headers and rows:
                        tables.append({
                            "page": page_num,
                            "index": t_idx,
                            "headers": headers,
                            "rows": rows[:200],  # limit rows
                            "row_count": len(rows),
                            "col_count": len(headers),
                        })
    except Exception as e:
        logger.warning(f"PDF table extraction failed: {e}")
    return tables


def _tables_to_markdown(tables: list[dict]) -> str:
    """Convert extracted tables to markdown format for LLM context."""
    if not tables:
        return ""
    parts = ["\n\n## Extracted Tables\n"]
    for t in tables[:20]:  # Limit to 20 tables
        parts.append(f"\n### Table (Page {t['page']})")
        # Header row
        parts.append("| " + " | ".join(t["headers"]) + " |")
        parts.append("| " + " | ".join("---" for _ in t["headers"]) + " |")
        # Data rows (limit to 50 for context)
        for row in t["rows"][:50]:
            # Pad or trim to match header count
            padded = (row + [""] * len(t["headers"]))[:len(t["headers"])]
            parts.append("| " + " | ".join(padded) + " |")
        if t["row_count"] > 50:
            parts.append(f"\n*...{t['row_count'] - 50} more rows*")
    return "\n".join(parts)


def _extract_text(file_path: str, file_type: str) -> str:
    try:
        if file_type == "pdf":
            import fitz
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
                if len(text) > MAX_TEXT_CHARS:
                    break
            doc.close()
            text = text[:MAX_TEXT_CHARS]

            # Also extract tables and append as markdown
            tables = _extract_tables_from_pdf(file_path)
            if tables:
                table_md = _tables_to_markdown(tables)
                text = text + table_md

            return text[:MAX_TEXT_CHARS]

        elif file_type == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                parts.append(f"## Sheet: {sheet}")
                rows_data = []
                headers = None
                for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                    cells = [str(c) if c is not None else "" for c in row]
                    if not any(cells):
                        continue
                    if row_idx == 0:
                        headers = cells
                        # Create markdown table header
                        parts.append("| " + " | ".join(cells) + " |")
                        parts.append("| " + " | ".join("---" for _ in cells) + " |")
                    else:
                        padded = (cells + [""] * len(headers or []))[:len(headers or cells)]
                        parts.append("| " + " | ".join(padded) + " |")
                    if len("\n".join(parts)) > MAX_TEXT_CHARS:
                        break
                if len("\n".join(parts)) > MAX_TEXT_CHARS:
                    break
            wb.close()
            return "\n".join(parts)[:MAX_TEXT_CHARS]

        elif file_type == "pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"## Slide {i}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                    # Extract tables from slides
                    if shape.has_table:
                        table = shape.table
                        headers = [cell.text.strip() for cell in table.rows[0].cells]
                        parts.append("| " + " | ".join(headers) + " |")
                        parts.append("| " + " | ".join("---" for _ in headers) + " |")
                        for row in table.rows[1:]:
                            cells = [cell.text.strip() for cell in row.cells]
                            parts.append("| " + " | ".join(cells) + " |")
                if len("\n".join(parts)) > MAX_TEXT_CHARS:
                    break
            return "\n".join(parts)[:MAX_TEXT_CHARS]

        elif file_type == "docx":
            import docx
            doc = docx.Document(file_path)
            parts = []
            # Extract paragraphs
            for p in doc.paragraphs:
                if p.text.strip():
                    parts.append(p.text.strip())
            # Extract tables
            for t_idx, table in enumerate(doc.tables):
                parts.append(f"\n### Table {t_idx + 1}")
                for r_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    if r_idx == 0:
                        parts.append("| " + " | ".join(cells) + " |")
                        parts.append("| " + " | ".join("---" for _ in cells) + " |")
                    else:
                        parts.append("| " + " | ".join(cells) + " |")
            return "\n".join(parts)[:MAX_TEXT_CHARS]

        elif file_type == "txt":
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read(MAX_TEXT_CHARS)

        return ""
    except Exception as e:
        logger.error(f"Text extraction failed for {file_path}: {e}")
        return ""


# ═══════════════════════════════════════════════════════════
# FOLDER CRUD
# ═══════════════════════════════════════════════════════════


async def create_folder(db: AsyncIOMotorDatabase, user_id: str, name: str, parent_id: str = None, color: str = None) -> dict:
    clean_name = name.strip()[:100]
    # Check for duplicate folder name in the same parent
    existing = await db.workspace_folders.find_one({
        "user_id": user_id, "parent_id": parent_id, "name": clean_name,
    })
    if existing:
        raise ValueError(f'A folder named "{clean_name}" already exists here')

    folder = {
        "user_id": user_id,
        "name": clean_name,
        "parent_id": parent_id,
        "color": color,
        "access_level": "workspace",
        "shared_with": [],
        "locked_by": None,
        "created_at": datetime.now(timezone.utc),
    }
    result = await db.workspace_folders.insert_one(folder)
    folder["_id"] = str(result.inserted_id)
    return folder


async def get_folder_contents(db: AsyncIOMotorDatabase, user_id: str, folder_id: str = None) -> dict:
    """Get folders and files at a given level. folder_id=None means root."""
    # Sub-folders
    folders = []
    async for f in db.workspace_folders.find(
        {"user_id": user_id, "parent_id": folder_id}
    ).sort("name", 1):
        # Count items inside this folder
        sub_folders = await db.workspace_folders.count_documents({"user_id": user_id, "parent_id": str(f["_id"])})
        sub_files = await db.workspace_files.count_documents({"user_id": user_id, "folder_id": str(f["_id"])})
        folders.append({
            "id": str(f["_id"]),
            "name": f["name"],
            "parent_id": f.get("parent_id"),
            "color": f.get("color"),
            "item_count": sub_folders + sub_files,
            "access_level": f.get("access_level", "workspace"),
            "shared_with": f.get("shared_with", []),
            "locked_by": f.get("locked_by"),
            "created_at": f.get("created_at"),
        })

    # Files at this level
    files = []
    async for f in db.workspace_files.find(
        {"user_id": user_id, "folder_id": folder_id},
        {"extracted_text": 0},
    ).sort("original_filename", 1):
        files.append({
            "id": str(f["_id"]),
            "original_filename": f["original_filename"],
            "file_type": f["file_type"],
            "char_count": f.get("char_count", 0),
            "page_count": f.get("page_count"),
            "status": f.get("status", "ready"),
            "folder_id": f.get("folder_id"),
            "access_level": f.get("access_level", "workspace"),
            "shared_with": f.get("shared_with", []),
            "locked_by": f.get("locked_by"),
            "created_at": f.get("created_at"),
        })

    return {"folders": folders, "files": files}


async def get_folder_path(db: AsyncIOMotorDatabase, folder_id: str) -> list[dict]:
    """Build breadcrumb path from root to this folder. Returns [{id, name}, ...]"""
    path = []
    current_id = folder_id
    seen = set()
    while current_id and current_id not in seen:
        seen.add(current_id)
        folder = await db.workspace_folders.find_one({"_id": ObjectId(current_id)})
        if not folder:
            break
        path.insert(0, {"id": str(folder["_id"]), "name": folder["name"]})
        current_id = folder.get("parent_id")
    return path


async def rename_folder(db: AsyncIOMotorDatabase, user_id: str, folder_id: str, name: str) -> bool:
    result = await db.workspace_folders.update_one(
        {"_id": ObjectId(folder_id), "user_id": user_id},
        {"$set": {"name": name.strip()[:100]}},
    )
    return result.modified_count > 0


async def move_folder(db: AsyncIOMotorDatabase, user_id: str, folder_id: str, new_parent_id: str = None) -> bool:
    # Prevent moving folder into itself or its descendants
    if new_parent_id:
        path = await get_folder_path(db, new_parent_id)
        if any(p["id"] == folder_id for p in path):
            return False  # Circular
    result = await db.workspace_folders.update_one(
        {"_id": ObjectId(folder_id), "user_id": user_id},
        {"$set": {"parent_id": new_parent_id}},
    )
    return result.modified_count > 0


async def delete_folder(db: AsyncIOMotorDatabase, user_id: str, folder_id: str) -> bool:
    """Cascade delete: remove folder, all subfolders, and all files within."""
    folder = await db.workspace_folders.find_one({"_id": ObjectId(folder_id), "user_id": user_id})
    if not folder:
        return False

    # Collect all descendant folder IDs
    to_delete = [folder_id]
    queue = [folder_id]
    while queue:
        current = queue.pop(0)
        async for sub in db.workspace_folders.find({"user_id": user_id, "parent_id": current}):
            sub_id = str(sub["_id"])
            to_delete.append(sub_id)
            queue.append(sub_id)

    # Delete files in all folders
    for fid in to_delete:
        async for f in db.workspace_files.find({"user_id": user_id, "folder_id": fid}):
            if os.path.exists(f.get("file_path", "")):
                os.remove(f["file_path"])
        await db.workspace_files.delete_many({"user_id": user_id, "folder_id": fid})

    # Delete all folders
    for fid in to_delete:
        await db.workspace_folders.delete_one({"_id": ObjectId(fid)})

    return True


async def move_file(db: AsyncIOMotorDatabase, user_id: str, file_id: str, folder_id: str = None) -> bool:
    result = await db.workspace_files.update_one(
        {"_id": ObjectId(file_id), "user_id": user_id},
        {"$set": {"folder_id": folder_id}},
    )
    return result.modified_count > 0


async def get_all_folders_tree(db: AsyncIOMotorDatabase, user_id: str) -> list[dict]:
    """Flat list of all folders for tree selector in move modal."""
    folders = []
    async for f in db.workspace_folders.find({"user_id": user_id}).sort("name", 1):
        folders.append({
            "id": str(f["_id"]),
            "name": f["name"],
            "parent_id": f.get("parent_id"),
        })
    return folders


# ═══════════════════════════════════════════════════════════
# FILE UPLOAD
# ═══════════════════════════════════════════════════════════


async def upload_workspace_file(
    db: AsyncIOMotorDatabase, user_id: str,
    file_content: bytes, original_filename: str,
    folder_id: str = None,
) -> dict:
    file_type = _detect_file_type(original_filename)
    if not file_type:
        raise ValueError(f"Unsupported file type. Allowed: {', '.join(ALLOWED_EXTENSIONS)}")

    # Check for duplicate filename in the same folder
    existing = await db.workspace_files.find_one({
        "user_id": user_id, "folder_id": folder_id, "original_filename": original_filename,
    })
    if existing:
        raise ValueError(f'A file named "{original_filename}" already exists in this folder')

    file_id = str(ObjectId())
    user_dir = os.path.join(settings.UPLOAD_DIR, user_id, "workspace")
    os.makedirs(user_dir, exist_ok=True)

    ext = os.path.splitext(original_filename)[1].lower()
    filename = f"{file_id}{ext}"
    file_path = os.path.join(user_dir, filename)

    with open(file_path, "wb") as f:
        f.write(file_content)

    extracted_text = _extract_text(file_path, file_type)
    char_count = len(extracted_text)

    # Extract structured tables for table intelligence (non-blocking — failures are silent)
    parsed_tables = []
    try:
        if file_type == "pdf":
            parsed_tables = _extract_tables_from_pdf(file_path)
        elif file_type == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            for sheet in wb.sheetnames[:10]:
                ws = wb[sheet]
                headers = None
                data_rows = []
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c or "").strip() for c in row]
                    if not any(cells):
                        continue
                    if headers is None:
                        headers = cells
                    else:
                        data_rows.append(cells)
                        row_count += 1
                        if row_count >= 200:
                            break
                if headers and data_rows:
                    parsed_tables.append({"page": 0, "index": 0, "headers": headers, "rows": data_rows, "row_count": len(data_rows), "col_count": len(headers), "sheet": sheet})
            wb.close()
    except Exception as e:
        logger.warning(f"Table extraction failed for {file_type}: {e}")
        parsed_tables = []

    page_count = None
    if file_type == "pdf":
        try:
            import fitz
            doc = fitz.open(file_path)
            page_count = len(doc)
            doc.close()
        except Exception:
            pass

    record = {
        "_id": ObjectId(file_id),
        "user_id": user_id,
        "workspace_id": user_id,
        "original_filename": original_filename,
        "file_type": file_type,
        "file_path": file_path,
        "extracted_text": extracted_text,
        "char_count": char_count,
        "page_count": page_count,
        "parsed_tables": parsed_tables[:50],  # max 50 tables per file
        "folder_id": folder_id,
        "access_level": "workspace",
        "shared_with": [],
        "locked_by": None,
        "status": "ready" if char_count > 0 else "failed",
        "error_message": None if char_count > 0 else "No text could be extracted",
        "created_at": datetime.now(timezone.utc),
    }
    await db.workspace_files.insert_one(record)
    record["_id"] = file_id
    return record


async def get_workspace_file(db: AsyncIOMotorDatabase, user_id: str, file_id: str) -> dict | None:
    f = await db.workspace_files.find_one({"_id": ObjectId(file_id), "user_id": user_id})
    if f:
        f["_id"] = str(f["_id"])
        f["id"] = f["_id"]
    return f


async def delete_workspace_file(db: AsyncIOMotorDatabase, user_id: str, file_id: str) -> bool:
    f = await db.workspace_files.find_one({"_id": ObjectId(file_id), "user_id": user_id})
    if not f:
        return False
    if os.path.exists(f["file_path"]):
        os.remove(f["file_path"])
    await db.workspace_files.delete_one({"_id": ObjectId(file_id)})
    return True


# ═══════════════════════════════════════════════════════════
# WORKSPACE CHAT (folder-aware)
# ═══════════════════════════════════════════════════════════


async def _collect_files_recursive(db: AsyncIOMotorDatabase, user_id: str, folder_id: str) -> list[dict]:
    """Collect all files in a folder and its subfolders recursively."""
    files = []
    async for f in db.workspace_files.find(
        {"user_id": user_id, "folder_id": folder_id, "status": "ready"},
        {"original_filename": 1, "file_type": 1, "extracted_text": 1, "char_count": 1, "folder_id": 1},
    ):
        files.append(f)

    # Recurse into subfolders
    async for sub in db.workspace_folders.find({"user_id": user_id, "parent_id": folder_id}):
        sub_files = await _collect_files_recursive(db, user_id, str(sub["_id"]))
        files.extend(sub_files)

    return files


async def chat_with_workspace(
    db: AsyncIOMotorDatabase, user: dict, message: str,
    conversation_id: str = None, file_id: str = None, folder_id: str = None,
):
    """Chat scoped to: a single file, a folder (recursive), or the entire workspace."""
    user_id = user["_id"]

    # Collect files based on scope
    if file_id:
        raw_files = [await db.workspace_files.find_one({"_id": ObjectId(file_id), "user_id": user_id})]
        raw_files = [f for f in raw_files if f]
        scope_label = "a single file"
    elif folder_id:
        raw_files = await _collect_files_recursive(db, user_id, folder_id)
        path = await get_folder_path(db, folder_id)
        scope_label = "/".join(p["name"] for p in path) if path else "folder"
    else:
        # Entire workspace — collect root files + all folder files
        raw_files = []
        async for f in db.workspace_files.find(
            {"user_id": user_id, "status": "ready"},
            {"original_filename": 1, "file_type": 1, "extracted_text": 1, "char_count": 1, "folder_id": 1},
        ).sort("char_count", -1).limit(30):
            raw_files.append(f)
        scope_label = "the entire workspace"

    if not raw_files:
        async def empty_gen():
            yield "I don't have any files in this scope yet. Upload some documents to get started!"
        return empty_gen(), None

    # Build file paths for context labels
    folder_path_cache = {}

    async def get_file_label(f):
        fid = f.get("folder_id")
        if fid and fid not in folder_path_cache:
            path = await get_folder_path(db, fid)
            folder_path_cache[fid] = "/".join(p["name"] for p in path)
        prefix = folder_path_cache.get(fid, "")
        name = f.get("original_filename", "unknown")
        return f"{prefix}/{name}" if prefix else name

    # Build context
    file_contexts = []
    total_chars = 0
    MAX_CONTEXT = 15000

    for f in raw_files[:20]:
        text = f.get("extracted_text", "")
        if not text:
            continue
        label = await get_file_label(f)
        snippet_len = min(3000, MAX_CONTEXT - total_chars) if total_chars < MAX_CONTEXT else 500
        snippet = text[:max(snippet_len, 500)]
        file_contexts.append(f"[File: {label}]\n{snippet}")
        total_chars += len(snippet)
        if total_chars >= MAX_CONTEXT:
            break

    context = "\n\n---\n\n".join(file_contexts)
    file_labels = [await get_file_label(f) for f in raw_files[:20]]

    system_prompt = f"""You are a workspace assistant with access to {len(raw_files)} files in {scope_label}.

RULES:
1. Answer using ONLY the workspace file content provided below.
2. ALWAYS cite sources in the format [File: folder/filename.ext] at the end of relevant sentences.
3. If information spans multiple files, cite each one.
4. If the answer is not in any file, say: "This information is not found in your workspace files."
5. Use markdown formatting.
6. Be concise but thorough.

Available files: {', '.join(file_labels[:15])}"""

    # Load conversation history
    history = []
    if conversation_id:
        conv = await db.workspace_conversations.find_one({"_id": ObjectId(conversation_id), "user_id": user_id})
        if conv:
            history = conv.get("messages", [])[-6:]

    messages = [{"role": "system", "content": system_prompt}]
    for msg in history:
        messages.append({"role": msg["role"], "content": msg["content"]})
    messages.append({"role": "user", "content": f"Question: {message}\n\nWorkspace Content:\n{context}"})

    response = await llm_service.completion(user=user, messages=messages, stream=True)

    # Save conversation
    if not conversation_id:
        conv_doc = {
            "user_id": user_id,
            "title": message[:80],
            "file_id": file_id,
            "folder_id": folder_id,
            "messages": [{"role": "user", "content": message}],
            "created_at": datetime.now(timezone.utc),
        }
        result = await db.workspace_conversations.insert_one(conv_doc)
        conversation_id = str(result.inserted_id)
    else:
        await db.workspace_conversations.update_one(
            {"_id": ObjectId(conversation_id)},
            {"$push": {"messages": {"role": "user", "content": message}}},
        )

    async def stream_gen():
        full_response = ""
        async for chunk in response:
            delta = chunk.choices[0].delta
            if delta.content:
                full_response += delta.content
                yield delta.content
        await db.workspace_conversations.update_one(
            {"_id": ObjectId(conversation_id)},
            {"$push": {"messages": {"role": "assistant", "content": full_response}}},
        )
        yield f"<!--conv_id:{conversation_id}-->"

    return stream_gen(), conversation_id
