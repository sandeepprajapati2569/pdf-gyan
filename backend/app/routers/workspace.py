"""
Workspace Router — folder-based file manager + chat.
"""

from fastapi import APIRouter, HTTPException, Depends, UploadFile, File, Query, Form
from fastapi.responses import StreamingResponse
from app.dependencies import get_current_user, get_user_db
from app.services.workspace_service import (
    create_folder, get_folder_contents, get_folder_path, rename_folder,
    move_folder, delete_folder, move_file, get_all_folders_tree,
    upload_workspace_file, get_workspace_file, delete_workspace_file,
    chat_with_workspace, set_access_level, get_access_info, ALLOWED_EXTENSIONS,
)
from pydantic import BaseModel
from typing import Optional
from bson import ObjectId

router = APIRouter(prefix="/api/workspace", tags=["workspace"])


# ═══ Folder CRUD ═══════════════════════════════════════════


class CreateFolderRequest(BaseModel):
    name: str
    parent_id: Optional[str] = None
    color: Optional[str] = None


class UpdateFolderRequest(BaseModel):
    name: Optional[str] = None
    parent_id: Optional[str] = "__unchanged__"  # sentinel to distinguish None from not-sent


class MoveFileRequest(BaseModel):
    folder_id: Optional[str] = None  # None = move to root


@router.post("/folders")
async def create_folder_endpoint(
    request: CreateFolderRequest,
    current_user: dict = Depends(get_current_user),
):
    db = await get_user_db(current_user)
    try:
        folder = await create_folder(db, current_user["_id"], request.name, request.parent_id, request.color)
    except ValueError as e:
        raise HTTPException(400, str(e))
    return folder


@router.get("/folders")
async def list_folder_contents(
    parent_id: Optional[str] = Query(None),
    current_user: dict = Depends(get_current_user),
):
    """Get folders and files at a given level. parent_id=null means root."""
    db = await get_user_db(current_user)
    contents = await get_folder_contents(db, current_user["_id"], parent_id)

    # Also return breadcrumb path if inside a folder
    path = []
    if parent_id:
        path = await get_folder_path(db, parent_id)

    return {**contents, "path": path, "current_folder_id": parent_id}


@router.get("/folders/tree")
async def get_folder_tree(current_user: dict = Depends(get_current_user)):
    """Flat list of all folders for the move-to-folder tree selector."""
    db = await get_user_db(current_user)
    return await get_all_folders_tree(db, current_user["_id"])


@router.get("/folders/{folder_id}")
async def get_folder_detail(
    folder_id: str,
    current_user: dict = Depends(get_current_user),
):
    db = await get_user_db(current_user)
    folder = await db.workspace_folders.find_one({"_id": ObjectId(folder_id), "user_id": current_user["_id"]})
    if not folder:
        raise HTTPException(404, "Folder not found")
    path = await get_folder_path(db, folder_id)
    return {"id": str(folder["_id"]), "name": folder["name"], "parent_id": folder.get("parent_id"), "color": folder.get("color"), "path": path}


@router.patch("/folders/{folder_id}")
async def update_folder_endpoint(
    folder_id: str,
    request: UpdateFolderRequest,
    current_user: dict = Depends(get_current_user),
):
    db = await get_user_db(current_user)
    if request.name:
        await rename_folder(db, current_user["_id"], folder_id, request.name)
    if request.parent_id != "__unchanged__":
        success = await move_folder(db, current_user["_id"], folder_id, request.parent_id)
        if not success:
            raise HTTPException(400, "Cannot move folder into itself or a subfolder")
    return {"status": "updated"}


@router.delete("/folders/{folder_id}")
async def delete_folder_endpoint(
    folder_id: str,
    current_user: dict = Depends(get_current_user),
):
    db = await get_user_db(current_user)
    deleted = await delete_folder(db, current_user["_id"], folder_id)
    if not deleted:
        raise HTTPException(404, "Folder not found")
    return {"status": "deleted"}


# ═══ File Management ═══════════════════════════════════════


@router.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    folder_id: Optional[str] = Form(None),
    current_user: dict = Depends(get_current_user),
):
    import os
    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(400, f"Unsupported file type '{ext}'. Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}")

    db = await get_user_db(current_user)
    content = await file.read()

    if not content:
        raise HTTPException(400, "Empty file")
    if len(content) > 50 * 1024 * 1024:
        raise HTTPException(400, "File too large. Maximum 50MB.")

    # Normalize folder_id (empty string → None)
    actual_folder_id = folder_id if folder_id and folder_id.strip() else None

    try:
        record = await upload_workspace_file(db, current_user["_id"], content, file.filename, actual_folder_id)
    except ValueError as e:
        raise HTTPException(400, str(e))
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"Upload failed: {str(e)}")

    return {
        "id": str(record["_id"]),
        "original_filename": record["original_filename"],
        "file_type": record["file_type"],
        "char_count": record["char_count"],
        "page_count": record.get("page_count"),
        "folder_id": record.get("folder_id"),
        "status": record["status"],
    }


# NOTE: Sub-path routes MUST come before the generic /files/{file_id} to avoid path conflicts

@router.get("/files/{file_id}/download")
async def download_file(file_id: str, current_user: dict = Depends(get_current_user)):
    """Serve the original file for download or in-browser viewing."""
    from fastapi.responses import Response
    import os
    db = await get_user_db(current_user)
    f = await get_workspace_file(db, current_user["_id"], file_id)
    if not f:
        raise HTTPException(404, "File not found")
    file_path = f.get("file_path", "")
    if not file_path or not os.path.exists(file_path):
        raise HTTPException(404, "File not found on disk")
    mime_map = {
        "pdf": "application/pdf",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "txt": "text/plain",
        "csv": "text/csv",
        "md": "text/markdown",
    }
    media_type = mime_map.get(f.get("file_type", ""), "application/octet-stream")
    filename = f.get("original_filename", "file")
    with open(file_path, "rb") as fh:
        content = fh.read()
    return Response(
        content=content,
        media_type=media_type,
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"',
            "Content-Length": str(len(content)),
        },
    )


@router.get("/files/{file_id}/preview")
async def preview_file(file_id: str, page: int = Query(1, ge=1), current_user: dict = Depends(get_current_user)):
    """Preview a workspace file. For PDFs returns a page image, for others returns extracted text."""
    from fastapi.responses import Response
    import os
    db = await get_user_db(current_user)
    f = await get_workspace_file(db, current_user["_id"], file_id)
    if not f:
        raise HTTPException(404, "File not found")

    file_path = f.get("file_path", "")
    file_type = f.get("file_type", "")

    if file_type == "pdf":
        if not file_path or not os.path.exists(file_path):
            raise HTTPException(400, f"PDF file not found on disk: {file_path}")
        try:
            import fitz
            pdf = fitz.open(file_path)
            total = len(pdf)
            if page < 1 or page > total:
                pdf.close()
                raise HTTPException(400, f"Page {page} out of range (1-{total})")
            pg = pdf[page - 1]
            pix = pg.get_pixmap(dpi=150)
            img_bytes = pix.tobytes("png")
            pdf.close()
            return Response(content=img_bytes, media_type="image/png", headers={"X-Total-Pages": str(total)})
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(500, f"PDF preview failed: {str(e)}")

    # PPTX: return structured slide data with layout info
    if file_type == "pptx" and file_path and os.path.exists(file_path):
        try:
            from pptx import Presentation
            from pptx.util import Emu
            prs = Presentation(file_path)
            slide_w = prs.slide_width or Emu(9144000)  # default 10"
            slide_h = prs.slide_height or Emu(6858000)  # default 7.5"
            w_px = slide_w / Emu(1)  # emu units
            h_px = slide_h / Emu(1)

            slides = []
            for i, slide in enumerate(prs.slides, 1):
                shapes = []
                for shape in slide.shapes:
                    s = {
                        "left": (shape.left or 0) / w_px * 100,    # percent of slide width
                        "top": (shape.top or 0) / h_px * 100,      # percent of slide height
                        "width": (shape.width or 0) / w_px * 100,
                        "height": (shape.height or 0) / h_px * 100,
                    }
                    if hasattr(shape, "text") and shape.text.strip():
                        # Extract paragraph-level formatting
                        paragraphs = []
                        if hasattr(shape, 'text_frame'):
                            for para in shape.text_frame.paragraphs:
                                p_text = para.text.strip()
                                if not p_text:
                                    continue
                                font_size = None
                                font_bold = False
                                font_color = None
                                align = str(para.alignment) if para.alignment else None
                                for run in para.runs:
                                    if run.font.size:
                                        font_size = round(run.font.size / Emu(12700))  # pt
                                    if run.font.bold:
                                        font_bold = True
                                    if run.font.color and run.font.color.rgb:
                                        font_color = f"#{run.font.color.rgb}"
                                paragraphs.append({
                                    "text": p_text,
                                    "font_size": font_size,
                                    "bold": font_bold,
                                    "color": font_color,
                                    "align": align,
                                })
                        s["type"] = "text"
                        s["paragraphs"] = paragraphs
                        shapes.append(s)

                    elif shape.has_table:
                        tbl = shape.table
                        headers = [cell.text.strip() for cell in tbl.rows[0].cells]
                        rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows[1:]]
                        s["type"] = "table"
                        s["headers"] = headers
                        s["rows"] = rows
                        shapes.append(s)

                # Try to get slide background color
                bg_color = None
                try:
                    bg = slide.background
                    if bg.fill and bg.fill.fore_color and bg.fill.fore_color.rgb:
                        bg_color = f"#{bg.fill.fore_color.rgb}"
                except Exception:
                    pass

                slides.append({"index": i, "shapes": shapes, "bg_color": bg_color})

            return {"file_type": "pptx", "slides": slides, "slide_count": len(slides), "aspect": round(w_px / h_px, 3)}
        except Exception:
            pass  # Fall through to text

    # For other types: return text content
    text = f.get("extracted_text", "") or ""
    return {"file_type": file_type, "text": text[:50000], "page_count": f.get("page_count")}


@router.get("/files/{file_id}/tables")
async def get_file_tables(file_id: str, current_user: dict = Depends(get_current_user)):
    """Get parsed tables from a workspace file."""
    db = await get_user_db(current_user)
    f = await get_workspace_file(db, current_user["_id"], file_id)
    if not f:
        raise HTTPException(404, "File not found")
    tables = f.get("parsed_tables", [])
    return {
        "file_id": file_id,
        "file_type": f.get("file_type"),
        "table_count": len(tables),
        "tables": tables,
    }


@router.get("/files/{file_id}/access")
async def get_file_access(file_id: str, current_user: dict = Depends(get_current_user)):
    db = await get_user_db(current_user)
    info = await get_access_info(db, current_user["_id"], "file", file_id)
    if not info:
        raise HTTPException(404, "File not found")
    return info


# Generic file GET must come AFTER all /files/{file_id}/subpath routes
@router.get("/files/{file_id}")
async def get_file(file_id: str, current_user: dict = Depends(get_current_user)):
    db = await get_user_db(current_user)
    f = await get_workspace_file(db, current_user["_id"], file_id)
    if not f:
        raise HTTPException(404, "File not found")
    return {
        "id": f["id"], "original_filename": f["original_filename"],
        "file_type": f["file_type"], "char_count": f.get("char_count", 0),
        "page_count": f.get("page_count"), "status": f.get("status"),
        "folder_id": f.get("folder_id"), "created_at": f.get("created_at"),
        "text_preview": (f.get("extracted_text") or "")[:500],
    }


@router.delete("/files/{file_id}")
async def remove_file(file_id: str, current_user: dict = Depends(get_current_user)):
    db = await get_user_db(current_user)
    if not await delete_workspace_file(db, current_user["_id"], file_id):
        raise HTTPException(404, "File not found")
    return {"status": "deleted"}


@router.post("/files/{file_id}/move")
async def move_file_endpoint(
    file_id: str,
    request: MoveFileRequest,
    current_user: dict = Depends(get_current_user),
):
    db = await get_user_db(current_user)
    if not await move_file(db, current_user["_id"], file_id, request.folder_id):
        raise HTTPException(404, "File not found")
    return {"status": "moved"}


# ═══ Chat ══════════════════════════════════════════════════


def _format_sse(chunk: str) -> str:
    normalized = chunk.replace("\r\n", "\n").replace("\r", "\n")
    return "".join(f"data: {line}\n" for line in normalized.split("\n")) + "\n"


class WorkspaceChatRequest(BaseModel):
    message: str
    conversation_id: Optional[str] = None
    file_id: Optional[str] = None
    folder_id: Optional[str] = None


@router.post("/chat")
async def workspace_chat(request: WorkspaceChatRequest, current_user: dict = Depends(get_current_user)):
    db = await get_user_db(current_user)
    try:
        gen, conv_id = await chat_with_workspace(
            db, current_user, request.message,
            request.conversation_id, request.file_id, request.folder_id,
        )
    except Exception as e:
        raise HTTPException(500, f"Chat failed: {str(e)}")

    async def stream():
        async for chunk in gen:
            yield _format_sse(chunk)
        yield "data: [DONE]\n\n"

    return StreamingResponse(stream(), media_type="text/event-stream", headers={"Cache-Control": "no-cache", "Connection": "keep-alive"})


@router.get("/chat/history")
async def workspace_chat_history(current_user: dict = Depends(get_current_user)):
    db = await get_user_db(current_user)
    cursor = db.workspace_conversations.find(
        {"user_id": current_user["_id"]},
        {"messages": {"$slice": -1}, "title": 1, "file_id": 1, "folder_id": 1, "created_at": 1},
    ).sort("created_at", -1).limit(50)
    return [{"id": str(c["_id"]), "title": c.get("title", "Untitled"), "file_id": c.get("file_id"), "folder_id": c.get("folder_id"), "created_at": c.get("created_at")} async for c in cursor]


@router.get("/chat/conversation/{conversation_id}")
async def get_workspace_conversation(conversation_id: str, current_user: dict = Depends(get_current_user)):
    db = await get_user_db(current_user)
    conv = await db.workspace_conversations.find_one({"_id": ObjectId(conversation_id), "user_id": current_user["_id"]})
    if not conv:
        raise HTTPException(404, "Conversation not found")
    return {"id": str(conv["_id"]), "title": conv.get("title"), "messages": conv.get("messages", []), "created_at": conv.get("created_at")}


# ═══════════════════════════════════════════════════════════
# ACCESS CONTROL
# ═══════════════════════════════════════════════════════════


class AccessRequest(BaseModel):
    access_level: str  # "workspace", "shared", "private"
    shared_with: Optional[list[str]] = None


@router.patch("/files/{file_id}/access")
async def set_file_access(file_id: str, request: AccessRequest, current_user: dict = Depends(get_current_user)):
    db = await get_user_db(current_user)
    result = await set_access_level(db, current_user["_id"], "file", file_id, request.access_level, request.shared_with)
    if not result:
        raise HTTPException(404, "File not found or invalid access level")
    return result


@router.patch("/folders/{folder_id}/access")
async def set_folder_access(folder_id: str, request: AccessRequest, current_user: dict = Depends(get_current_user)):
    db = await get_user_db(current_user)
    result = await set_access_level(db, current_user["_id"], "folder", folder_id, request.access_level, request.shared_with)
    if not result:
        raise HTTPException(404, "Folder not found or invalid access level")
    return result


@router.get("/folders/{folder_id}/access")
async def get_folder_access(folder_id: str, current_user: dict = Depends(get_current_user)):
    db = await get_user_db(current_user)
    info = await get_access_info(db, current_user["_id"], "folder", folder_id)
    if not info:
        raise HTTPException(404, "Folder not found")
    return info
